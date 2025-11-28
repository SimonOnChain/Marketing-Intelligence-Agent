"""Create PowerPoint presentation for Marketing Intelligence Agent."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BLUE = RgbColor(0x1a, 0x1a, 0x2e)
ACCENT_BLUE = RgbColor(0x00, 0xd4, 0xff)
WHITE = RgbColor(0xff, 0xff, 0xff)
GRAY = RgbColor(0xaa, 0xaa, 0xaa)

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bullets, code_snippet=None):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = RgbColor(0x16, 0x21, 0x3e)
    title_bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Content area
    if code_snippet:
        # Split layout - bullets on left, code on right
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(5.5))
    else:
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5.5))

    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(24)
        p.font.color.rgb = WHITE
        p.space_before = Pt(12)

    # Code snippet if provided
    if code_snippet:
        code_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(1.5), Inches(6.3), Inches(5.5))
        code_box.fill.solid()
        code_box.fill.fore_color.rgb = RgbColor(0x0d, 0x11, 0x17)
        code_box.line.color.rgb = RgbColor(0x30, 0x36, 0x3d)

        code_text = slide.shapes.add_textbox(Inches(6.7), Inches(1.7), Inches(5.9), Inches(5.1))
        tf = code_text.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = code_snippet
        p.font.size = Pt(12)
        p.font.name = "Consolas"
        p.font.color.rgb = RgbColor(0x9c, 0xdc, 0xfe)

    return slide

def add_architecture_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Systemarchitektur - Multi-Agent Pipeline"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Flow boxes
    boxes = [
        ("User Query", 0.5, 1.8),
        ("Orchestrator\n(Intent Classification)", 2.5, 1.8),
        ("Sales Agent", 5.5, 1.3),
        ("Sentiment Agent\n(RAG)", 5.5, 2.8),
        ("Forecast Agent\n(Prophet ML)", 5.5, 4.3),
        ("Synthesizer", 8.5, 2.8),
        ("Response", 11, 2.8),
    ]

    for text, x, y in boxes:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.8), Inches(0.9))
        box.fill.solid()
        if "Agent" in text or "Orchestrator" in text or "Synthesizer" in text:
            box.fill.fore_color.rgb = RgbColor(0x00, 0x7a, 0xcc)
        else:
            box.fill.fore_color.rgb = RgbColor(0x2d, 0x3a, 0x4a)
        box.line.color.rgb = ACCENT_BLUE

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Description
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12), Inches(1.5))
    tf = desc_box.text_frame
    p = tf.paragraphs[0]
    p.text = "LangGraph State Machine: Orchestriert spezialisierte Agenten basierend auf Intent-Klassifikation"
    p.font.size = Pt(20)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

    return slide

# ============ CREATE SLIDES ============

# Slide 1: Title
add_title_slide(
    prs,
    "Marketing Intelligence Agent",
    "Data Science, Machine Learning & AI\nSimon Jokani"
)

# Slide 2: Agenda
add_content_slide(prs, "Agenda", [
    "Projektübersicht & Problemstellung",
    "Multi-Agent Architektur (LangGraph)",
    "RAG Pipeline - Retrieval-Augmented Generation",
    "Prophet ML - Zeitreihen-Forecasting",
    "Hybrid Search (Vektor + Lexikalisch)",
    "Live Demo & Ergebnisse"
])

# Slide 3: Problem & Solution
add_content_slide(prs, "Problemstellung & Lösung", [
    "Problem: Marketing-Analysten verbringen Stunden mit manuellen Datenabfragen",
    "Datenbasis: Olist E-Commerce Dataset (100K+ Orders, 40K+ Reviews)",
    "Lösung: KI-gestützter Analyst für natürlichsprachliche Fragen",
    "Drei spezialisierte Agenten: Sales, Sentiment, Forecast",
    "Automatische Intent-Erkennung und Agent-Routing"
])

# Slide 4: Architecture
add_architecture_slide(prs)

# Slide 5: LangGraph
add_content_slide(prs, "LangGraph - Multi-Agent Orchestrierung", [
    "State Machine Pattern für Agent-Koordination",
    "TypedDict für typsichere Zustandsverwaltung",
    "Conditional Edges für dynamisches Routing",
    "Parallele Agent-Ausführung möglich"
], code_snippet="""# LangGraph State Machine
class AgentState(TypedDict):
    query: str
    intent: str
    agent_outputs: Dict
    final_response: str

graph = StateGraph(AgentState)
graph.add_node("classify", classify_intent)
graph.add_node("sales", sales_agent)
graph.add_node("sentiment", sentiment_agent)
graph.add_node("forecast", forecast_agent)
graph.add_conditional_edges(
    "classify",
    route_to_agent
)""")

# Slide 6: RAG Pipeline
add_content_slide(prs, "RAG - Retrieval-Augmented Generation", [
    "Embedding Model: sentence-transformers/MiniLM",
    "Vektor-Datenbank: Qdrant Cloud",
    "40K+ Reviews als Wissensbasis indexiert",
    "Semantische Suche für relevante Kontexte",
    "LLM generiert Antwort basierend auf Retrieval"
], code_snippet="""# RAG Pipeline
embeddings = SentenceTransformer(
    'all-MiniLM-L6-v2'
)

# Semantic search
results = qdrant_client.search(
    collection="reviews",
    query_vector=embeddings.encode(query),
    limit=10
)

# Generate with context
response = llm.generate(
    context=results,
    query=user_query
)""")

# Slide 7: Hybrid Search
add_content_slide(prs, "Hybrid Search - Vektor + Lexikalisch", [
    "Kombination aus semantischer und keyword-basierter Suche",
    "Vektor-Suche: Findet semantisch ähnliche Inhalte",
    "BM25 Lexikalische Suche: Exakte Keyword-Matches",
    "Reciprocal Rank Fusion für Score-Kombination",
    "Bessere Recall als reine Vektor-Suche"
], code_snippet="""# Hybrid Search
def hybrid_search(query, k=10):
    # Vector search
    vec_results = vector_search(
        query, k=k*2
    )

    # Lexical BM25 search
    lex_results = bm25_search(
        query, k=k*2
    )

    # Reciprocal Rank Fusion
    return rrf_combine(
        vec_results,
        lex_results,
        k=k
    )""")

# Slide 8: Prophet Forecasting
add_content_slide(prs, "Prophet ML - Zeitreihen-Forecasting", [
    "Facebook Prophet für Revenue-Prognosen",
    "Automatische Trend- und Saisonalitätserkennung",
    "Log-Transformation für stabile Vorhersagen",
    "8-Wochen Forecast mit Konfidenzintervall",
    "Handling von Feiertagen und Anomalien"
], code_snippet="""# Prophet Forecasting
from prophet import Prophet

# Prepare data
df = pd.DataFrame({
    'ds': dates,
    'y': np.log1p(revenue)  # Log transform
})

# Fit model
model = Prophet(
    yearly_seasonality=True,
    weekly_seasonality=True
)
model.fit(df)

# Predict
future = model.make_future_dataframe(
    periods=8, freq='W'
)
forecast = model.predict(future)
forecast['yhat'] = np.expm1(forecast['yhat'])""")

# Slide 9: Tech Stack
add_content_slide(prs, "Technologie-Stack", [
    "LLMs: xAI Grok, Groq (Llama), AWS Bedrock",
    "Embeddings: sentence-transformers MiniLM",
    "Vector DB: Qdrant Cloud",
    "ML: Prophet, scikit-learn, PyTorch",
    "Orchestration: LangGraph, LangChain",
    "Monitoring: Langfuse (Tracing & Analytics)",
    "Backend: FastAPI, Uvicorn"
])

# Slide 10: Results & Demo
add_content_slide(prs, "Ergebnisse & Live Demo", [
    "Automatische Intent-Klassifikation: >95% Accuracy",
    "RAG-basierte Sentiment-Analyse auf 40K Reviews",
    "Prophet Forecast: +9% Wachstumsprognose",
    "Antwortzeit: ~3-5 Sekunden pro Query",
    "Deployed auf AWS EC2 (weltweit erreichbar)"
])

# Slide 11: Conclusion
add_content_slide(prs, "Zusammenfassung & Learnings", [
    "Multi-Agent Systeme ermöglichen spezialisierte KI-Lösungen",
    "RAG verbessert LLM-Antworten mit domänenspezifischem Wissen",
    "Hybrid Search kombiniert Stärken beider Ansätze",
    "Prophet liefert interpretierbare ML-Forecasts",
    "LangGraph vereinfacht komplexe Agent-Workflows"
])

# Slide 12: Q&A
add_title_slide(
    prs,
    "Fragen?",
    "Live Demo: http://3.121.239.209:8501\nGitHub: github.com/SimonOnChain/Marketing-Intelligence-Agent"
)

# Save
output_path = Path("C:/Users/Simon/Project_Wifi/presentation")
output_path.mkdir(exist_ok=True)
prs.save(output_path / "Marketing_Intelligence_Agent_Presentation.pptx")
print(f"Presentation saved to: {output_path / 'Marketing_Intelligence_Agent_Presentation.pptx'}")
